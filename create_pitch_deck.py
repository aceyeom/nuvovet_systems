#!/usr/bin/env python3
"""
NuvoVet Pitch Deck Generator
Professional, dynamic pitch deck for veterinary DUR system
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import math

# ── Color Palette ──
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFB)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x8D)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_AMBER = RGBColor(0xFF, 0xB0, 0x3B)
MID_NAVY = RGBColor(0x2D, 0x3A, 0x8C)
SOFT_NAVY = RGBColor(0x3A, 0x4A, 0xA8)
LIGHT_ICE = RGBColor(0xE8, 0xEF, 0xFD)
CHARCOAL = RGBColor(0x33, 0x33, 0x44)
TRANSPARENT_NAVY = RGBColor(0x15, 0x1D, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def emu(inches):
    return Inches(inches)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    if line_width:
        shape.line.width = Pt(line_width)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri', space_before=0, space_after=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def dark_bg(slide):
    """Full dark navy background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def light_bg(slide):
    """Light background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

def add_nav_dots(slide, current_index, total=12):
    """Add navigation dots at bottom"""
    dot_size = Emu(Inches(0.12))
    spacing = Emu(Inches(0.25))
    total_width = total * Inches(0.25)
    start_x = Emu((W - Emu(total_width)) // 2)
    y = Emu(H - Inches(0.4))
    for i in range(total):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, start_x + Emu(i * Inches(0.25)), y, dot_size, dot_size)
        dot.line.fill.background()
        dot.fill.solid()
        if i == current_index:
            dot.fill.fore_color.rgb = ICE_BLUE if current_index in [0, 11] else NAVY
        else:
            dot.fill.fore_color.rgb = RGBColor(0x55, 0x60, 0x90) if current_index in [0, 11] else RGBColor(0xCC, 0xCC, 0xDD)

def add_section_icon(slide, left, top, size=0.5, icon_text="", bg_color=NAVY, text_color=WHITE):
    """Small circle with icon text inside"""
    circle = add_shape(slide, MSO_SHAPE.OVAL, emu(left), emu(top), emu(size), emu(size), fill_color=bg_color)
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size * 20))
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return circle

def add_stat_callout(slide, left, top, number, label, number_color=NAVY, bg_color=None, width=2.5, height=1.5):
    """Large number stat callout"""
    if bg_color:
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height), fill_color=bg_color)
        card.line.fill.background()
        # Set corner radius
        try:
            card._element.attrib['adjLst'] = ''
        except:
            pass

    txBox = slide.shapes.add_textbox(emu(left), emu(top + 0.1), emu(width), emu(height * 0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = number_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(emu(left), emu(top + height * 0.55), emu(width), emu(height * 0.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(13)
    p2.font.bold = False
    p2.font.color.rgb = CHARCOAL
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER

def add_card(slide, left, top, width, height, fill=WHITE, shadow=True, corner_radius=None):
    """Add a rounded rectangle card"""
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height), fill_color=fill)
    card.line.fill.background()
    return card

def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """Simulated gradient with overlapping shapes"""
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(width), emu(height), fill_color=color1)
    bar.line.fill.background()
    return bar


# ============================================================
# SLIDE 1: COVER (Dark Background)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
dark_bg(slide)

# Decorative geometric elements - large angled shape
pts = [
    (Inches(0), Inches(0)),
    (Inches(6), Inches(0)),
    (Inches(4.5), Inches(7.5)),
    (Inches(0), Inches(7.5)),
]
# Left decorative panel
left_panel = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.5), H, fill_color=TRANSPARENT_NAVY)
left_panel.line.fill.background()
# Make it semi-transparent via a lighter navy
left_panel.fill.fore_color.rgb = RGBColor(0x17, 0x20, 0x50)

# Diagonal accent line
diag = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(0), Inches(0.08), H, fill_color=ICE_BLUE)
diag.line.fill.background()
diag.rotation = 3.0

# Large decorative circles (abstract/molecular feel)
for cx, cy, sz, alpha_color in [
    (10.5, 1.0, 2.0, RGBColor(0x25, 0x30, 0x70)),
    (11.5, 5.5, 1.5, RGBColor(0x22, 0x2D, 0x68)),
    (9.0, 6.0, 1.0, RGBColor(0x20, 0x2A, 0x65)),
]:
    c = add_shape(slide, MSO_SHAPE.OVAL, emu(cx), emu(cy), emu(sz), emu(sz), fill_color=alpha_color)
    c.line.fill.background()

# Small dot grid pattern on the right side
for row in range(8):
    for col in range(6):
        x = 8.0 + col * 0.35
        y = 0.5 + row * 0.35
        dot = add_shape(slide, MSO_SHAPE.OVAL, emu(x), emu(y), emu(0.06), emu(0.06), fill_color=RGBColor(0x35, 0x42, 0x85))
        dot.line.fill.background()

# Ice blue accent bar at top
accent_top = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.6), Inches(0.06), fill_color=ICE_BLUE)
accent_top.line.fill.background()

# Logo / Brand
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.6), "NUVOVET",
             font_size=16, bold=True, color=ICE_BLUE, font_name='Calibri')

# Main title
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(4.2), Inches(1.2),
             "Veterinary Drug\nUtilization Review",
             font_size=42, bold=True, color=WHITE, font_name='Calibri')

# Subtitle
add_text_box(slide, Inches(0.8), Inches(3.5), Inches(4.2), Inches(0.8),
             "동물 의료의 처방 안전망",
             font_size=22, bold=False, color=ICE_BLUE, font_name='Calibri')

# Description
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(4.0), Inches(1.2),
             "국내 유일의 수의학 약물 안전성 검증 시스템\n5,600+ 동물병원을 위한 처방 보조 인프라",
             font_size=14, bold=False, color=RGBColor(0x99, 0xAA, 0xCC), font_name='Calibri')

# Key stats on the right side - glass-card style
for i, (num, label) in enumerate([
    ("641+", "Drug Database"),
    ("10", "Safety Rules"),
    ("<1s", "Analysis Time"),
]):
    y = 2.5 + i * 1.4
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(6.5), emu(y), emu(3.0), emu(1.1), fill_color=RGBColor(0x22, 0x2E, 0x6A))
    card.line.fill.background()
    # Add a left accent bar on card
    accent = add_shape(slide, MSO_SHAPE.RECTANGLE, emu(6.5), emu(y), emu(0.06), emu(1.1), fill_color=ICE_BLUE)
    accent.line.fill.background()

    add_text_box(slide, emu(6.85), emu(y + 0.12), emu(2.5), emu(0.55), num,
                 font_size=30, bold=True, color=WHITE, font_name='Calibri')
    add_text_box(slide, emu(6.85), emu(y + 0.65), emu(2.5), emu(0.35), label,
                 font_size=12, bold=False, color=ICE_BLUE, font_name='Calibri')

# Bottom tagline
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.5),
             "말 못하는 생명을 위한, 데이터 기반 처방 안전망",
             font_size=13, bold=False, color=RGBColor(0x88, 0x99, 0xBB), font_name='Calibri')

add_nav_dots(slide, 0, 12)


# ============================================================
# SLIDE 2: THE PROBLEM (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

# Top navy header strip
header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "THE PROBLEM",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "왜 동물 의료에 DUR이 필요한가?",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Section label icon
add_section_icon(slide, 12.2, 0.2, 0.5, "!", ACCENT_RED, WHITE)

# Three problem cards in a row
problems = [
    ("70%", "오프라벨 처방",
     "국내 동물병원 항생제의 약 70%가\n인체용 의약품으로 처방됩니다.\n대사 능력이 다른 동물에게\n부작용 검증 없이 투약하고 있습니다.",
     ACCENT_RED),
    ("0", "DUR 시스템 전무",
     "미국은 Instinct Science 등\n처방 보조가 필수 인프라이나,\n국내에는 약리학적 교차 검증을\n수행하는 시스템이 전무합니다.",
     ACCENT_AMBER),
    ("~50%", "부작용/오진 불만",
     "최근 3년 소비자연맹 통계 상\n소비자 불만의 절반이\n'치료 부작용 및 오진'에서\n비롯되었습니다.",
     ACCENT_RED),
]

for i, (stat, title, desc, accent) in enumerate(problems):
    x = 0.6 + i * 4.2
    y = 1.4

    # Card
    card = add_card(slide, x, y, 3.8, 4.8, WHITE)

    # Top accent bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(3.8), emu(0.06), fill_color=accent)

    # Large stat number
    add_text_box(slide, emu(x + 0.4), emu(y + 0.4), emu(3.0), emu(0.9), stat,
                 font_size=48, bold=True, color=accent, font_name='Calibri', alignment=PP_ALIGN.LEFT)

    # Title
    add_text_box(slide, emu(x + 0.4), emu(y + 1.3), emu(3.0), emu(0.45), title,
                 font_size=18, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.LEFT)

    # Divider line
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 0.4), emu(y + 1.85), emu(1.5), emu(0.03), fill_color=ICE_BLUE)

    # Description
    add_text_box(slide, emu(x + 0.4), emu(y + 2.1), emu(3.0), emu(2.5), desc,
                 font_size=13, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.LEFT)

# Bottom impact statement
impact_bar = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(0.6), emu(6.5), emu(12.1), emu(0.7), fill_color=NAVY)
impact_bar.line.fill.background()
add_text_box(slide, emu(1.0), emu(6.55), emu(11.5), emu(0.6),
             "수의사의 용량 계산 실수 하나가 동물의 생명을 앗아가고, 병원에는 의료 소송 리스크로 돌아옵니다.",
             font_size=14, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 1, 12)


# ============================================================
# SLIDE 3: THE SOLUTION (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

# Top navy header strip
header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "THE SOLUTION",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "NuvoVet: 한국형 동물 맞춤형 DUR 시스템",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

add_section_icon(slide, 12.2, 0.2, 0.5, "V", ACCENT_GREEN, WHITE)

# Left column - main value prop
add_text_box(slide, emu(0.8), emu(1.3), emu(5.5), emu(0.6),
             "처방 패러다임의 혁신",
             font_size=26, bold=True, color=NAVY, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(1.9), emu(5.5), emu(1.8),
             "수의사 개인의 경험 중심에서\n데이터 중심으로 전환하여\n완벽한 투약 안전망을 구축합니다.\n\n기존 EMR 시스템과 API로 연동되어\n처방 단계에서 자동으로 안전성을 검증합니다.",
             font_size=14, bold=False, color=CHARCOAL, font_name='Calibri')

# Solution feature cards (2x2 grid on right)
features = [
    ("Drug-Drug Interaction", "약물 간 상호작용을\n심각도별 실시간 분석", "DDI"),
    ("Organ Safety Scoring", "간/신장 누적 부담을\n정량적 점수로 산출", "ORG"),
    ("Breed Sensitivity", "품종별 유전적 취약점\n자동 감지 및 차단", "DNA"),
    ("Auto Dose Correction", "체중/종/상호작용 기반\n투여량 자동 보정", "Rx"),
]

for i, (title, desc, icon) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 6.8 + col * 3.1
    y = 1.3 + row * 2.1

    card = add_card(slide, x, y, 2.8, 1.8, WHITE)

    # Icon circle
    add_section_icon(slide, x + 0.25, y + 0.25, 0.55, icon, NAVY, WHITE)

    add_text_box(slide, emu(x + 1.0), emu(y + 0.25), emu(1.6), emu(0.4), title,
                 font_size=11, bold=True, color=NAVY, font_name='Calibri')
    add_text_box(slide, emu(x + 0.25), emu(y + 0.9), emu(2.3), emu(0.8), desc,
                 font_size=11, bold=False, color=CHARCOAL, font_name='Calibri')

# Bottom differentiator bar
diff_bar = add_card(slide, 0.6, 5.8, 12.1, 1.3, NAVY)
items = [
    ("EMR 연동", "API 기반 제로 프릭션"),
    ("다차원 검증", "병력 + 알레르기 + 품종 교차"),
    ("하드스탑 기능", "치명적 처방 원천 차단"),
    ("빅데이터 플랫폼", "처방기록 → 수의학 연구"),
]
for i, (title, sub) in enumerate(items):
    x = 1.0 + i * 3.0
    add_text_box(slide, emu(x), emu(5.95), emu(2.5), emu(0.35), title,
                 font_size=14, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, emu(x), emu(6.3), emu(2.5), emu(0.35), sub,
                 font_size=10, bold=False, color=ICE_BLUE, font_name='Calibri', alignment=PP_ALIGN.CENTER)
    if i < 3:
        # Vertical separator
        add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 2.7), emu(6.0), emu(0.02), emu(0.8), fill_color=RGBColor(0x35, 0x42, 0x85))

add_nav_dots(slide, 2, 12)


# ============================================================
# SLIDE 4: CORE FEATURES (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "CORE FEATURES",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "8대 안전 검증 엔진",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# 8 feature cards in 2 rows of 4
safety_rules = [
    ("01", "약물 상호작용", "Drug-Drug Interaction", "심각도별 분석\n(Critical/Moderate/Minor)", ACCENT_RED),
    ("02", "장기 부담 산출", "Organ Burden Score", "간/신장 누적 부담\n정량적 점수화", ACCENT_AMBER),
    ("03", "종 특이 독성", "Species Toxicity", "고양이 아세트아미노펜\n등 치명적 약물 차단", ACCENT_RED),
    ("04", "품종 유전 감수성", "Breed Genetics", "MDR1 변이(콜리 등)\n자동 감지/하드스탑", ACCENT_AMBER),
    ("05", "용량 자동 보정", "Auto Dose Calc", "체중/종/상호작용 기반\n투여량 역산", NAVY),
    ("06", "CYP 효소 분석", "CYP Metabolism", "CYP3A4/2D6 억제\n유도 영향 분석", NAVY),
    ("07", "QT 연장 위험", "QT Prolongation", "치명적 부정맥\n(토르사드) 검증", ACCENT_RED),
    ("08", "출혈 위험 합산", "Bleeding Risk", "항응고 + NSAID 등\n출혈 위험 합산 경고", ACCENT_AMBER),
]

for i, (num, title_ko, title_en, desc, accent) in enumerate(safety_rules):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.15
    y = 1.2 + row * 3.0

    card = add_card(slide, x, y, 2.9, 2.7, WHITE)

    # Number badge
    badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + 0.2), emu(y + 0.2), emu(0.55), emu(0.45), fill_color=accent)
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Korean title
    add_text_box(slide, emu(x + 0.9), emu(y + 0.2), emu(1.8), emu(0.45), title_ko,
                 font_size=14, bold=True, color=NAVY, font_name='Calibri')
    # English subtitle
    add_text_box(slide, emu(x + 0.2), emu(y + 0.75), emu(2.5), emu(0.3), title_en,
                 font_size=10, bold=False, color=SOFT_NAVY, font_name='Calibri')
    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 0.2), emu(y + 1.1), emu(1.2), emu(0.025), fill_color=ICE_BLUE)
    # Description
    add_text_box(slide, emu(x + 0.2), emu(y + 1.25), emu(2.5), emu(1.3), desc,
                 font_size=11, bold=False, color=CHARCOAL, font_name='Calibri')

add_nav_dots(slide, 3, 12)


# ============================================================
# SLIDE 5: HOW IT WORKS (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "HOW IT WORKS",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "1초 만에 처방전 안전성 검증",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Process flow - horizontal steps with arrows
steps = [
    ("1", "EMR 입력", "처방전 자동 수신\n또는 OCR 스캔", "EMR"),
    ("2", "약물 식별", "641+ 약물 DB에서\n성분/분류 매칭", "DB"),
    ("3", "다차원 분석", "10대 안전 규칙\n교차 검증 수행", "DUR"),
    ("4", "결과 리포트", "심각도별 경고 +\n대안 처방 제시", "Rx"),
]

# Connecting line
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(2.0), emu(3.15), emu(9.5), emu(0.04), fill_color=ICE_BLUE)

for i, (num, title, desc, icon) in enumerate(steps):
    x = 0.8 + i * 3.15

    # Circle step indicator
    circle = add_shape(slide, MSO_SHAPE.OVAL, emu(x + 0.9), emu(2.3), emu(0.8), emu(0.8), fill_color=NAVY)
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Step number
    add_text_box(slide, emu(x), emu(1.5), emu(2.8), emu(0.4), f"STEP {num}",
                 font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

    # Card below
    card = add_card(slide, x, 3.5, 2.8, 2.2, WHITE)
    add_text_box(slide, emu(x + 0.3), emu(3.7), emu(2.2), emu(0.4), title,
                 font_size=16, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 0.5), emu(4.15), emu(1.8), emu(0.025), fill_color=ICE_BLUE)
    add_text_box(slide, emu(x + 0.3), emu(4.35), emu(2.2), emu(1.2), desc,
                 font_size=12, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Arrow indicators between steps
for i in range(3):
    x = 3.2 + i * 3.15
    add_text_box(slide, emu(x), emu(2.4), emu(0.5), emu(0.5), "→",
                 font_size=24, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Bottom feature bar
feature_bar = add_card(slide, 0.6, 6.0, 12.1, 1.1, WHITE)
techs = ["React + FastAPI", "641+ JSONL Drug Records", "10 DUR Safety Rules", "Real-time OCR via Claude Vision"]
for i, t in enumerate(techs):
    x = 0.8 + i * 3.05
    dot = add_shape(slide, MSO_SHAPE.OVAL, emu(x), emu(6.35), emu(0.15), emu(0.15), fill_color=NAVY)
    dot.line.fill.background()
    add_text_box(slide, emu(x + 0.25), emu(6.25), emu(2.8), emu(0.4), t,
                 font_size=11, bold=False, color=NAVY, font_name='Calibri')

add_nav_dots(slide, 4, 12)


# ============================================================
# SLIDE 6: ORGAN SAFETY VISUALIZATION (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "ORGAN SAFETY",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "5대 장기 부담 시각화 시스템",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Left side - anatomy overview card
anatomy_card = add_card(slide, 0.6, 1.3, 5.5, 5.6, WHITE)

# Simulated anatomy diagram - body outline
body = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(1.5), emu(1.8), emu(3.5), emu(4.5), fill_color=LIGHT_ICE)
body.line.fill.background()

# Organ indicators on body
organs = [
    ("Brain", "뇌", 2.8, 2.1, 0.7, RGBColor(0x9B, 0x59, 0xB6)),
    ("Heart", "심장", 2.3, 3.1, 0.65, ACCENT_RED),
    ("Liver", "간", 3.4, 3.1, 0.65, ACCENT_AMBER),
    ("Blood", "혈액", 2.8, 4.0, 0.65, RGBColor(0xE7, 0x4C, 0x3C)),
    ("Kidney", "신장", 2.8, 4.8, 0.65, RGBColor(0x27, 0xAE, 0x60)),
]

for name, name_ko, ox, oy, sz, color in organs:
    organ_circle = add_shape(slide, MSO_SHAPE.OVAL, emu(ox), emu(oy), emu(sz), emu(sz), fill_color=color)
    organ_circle.line.fill.background()
    tf = organ_circle.text_frame
    p = tf.paragraphs[0]
    p.text = name_ko
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Right side - organ burden cards
for i, (name, name_ko, _, _, _, color) in enumerate(organs):
    x = 6.5
    y = 1.3 + i * 1.1

    row_card = add_card(slide, x, y, 6.2, 0.95, WHITE)

    # Color indicator
    add_shape(slide, MSO_SHAPE.OVAL, emu(x + 0.2), emu(y + 0.2), emu(0.55), emu(0.55), fill_color=color)

    # Organ name
    add_text_box(slide, emu(x + 0.9), emu(y + 0.1), emu(1.2), emu(0.35), name_ko,
                 font_size=14, bold=True, color=NAVY, font_name='Calibri')
    add_text_box(slide, emu(x + 0.9), emu(y + 0.45), emu(1.2), emu(0.3), name,
                 font_size=9, bold=False, color=CHARCOAL, font_name='Calibri')

    # Simulated burden bar
    bar_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + 2.3), emu(y + 0.35), emu(3.0), emu(0.25), fill_color=LIGHT_ICE)
    bar_bg.line.fill.background()
    bar_widths = [1.8, 2.4, 1.2, 2.1, 0.9]
    bar_fill = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + 2.3), emu(y + 0.35), emu(bar_widths[i]), emu(0.25), fill_color=color)
    bar_fill.line.fill.background()

    # Score
    scores = ["60", "80", "40", "70", "30"]
    add_text_box(slide, emu(x + 5.5), emu(y + 0.2), emu(0.6), emu(0.5), scores[i],
                 font_size=18, bold=True, color=color, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Bottom description
add_text_box(slide, emu(0.8), emu(6.95), emu(12), emu(0.4),
             "각 약물이 장기에 가하는 부담을 0-100 점수로 정량화  |  다약제 처방 시 누적 부담 자동 합산  |  종별(개/고양이) 차등 적용",
             font_size=11, bold=False, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 5, 12)


# ============================================================
# SLIDE 7: MARKET OPPORTUNITY (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "MARKET OPPORTUNITY",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "시장 분석 및 규모 (TAM / SAM / SOM)",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# TAM / SAM / SOM concentric circles visualization
# TAM - largest circle
tam = add_shape(slide, MSO_SHAPE.OVAL, emu(0.8), emu(1.3), emu(5.5), emu(5.5), fill_color=LIGHT_ICE)
tam.line.color.rgb = ICE_BLUE
tam.line.width = Pt(2)
tam.line.fill.solid()

# SAM - medium circle
sam = add_shape(slide, MSO_SHAPE.OVAL, emu(1.55), emu(2.05), emu(4.0), emu(4.0), fill_color=RGBColor(0xB8, 0xCD, 0xF0))
sam.line.color.rgb = NAVY
sam.line.width = Pt(1.5)
sam.line.fill.solid()

# SOM - small circle
som = add_shape(slide, MSO_SHAPE.OVAL, emu(2.3), emu(2.8), emu(2.5), emu(2.5), fill_color=NAVY)
som.line.fill.background()

# Labels
add_text_box(slide, emu(2.55), emu(3.5), emu(2.0), emu(0.7), "SOM\n₩3.4B",
             font_size=16, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)
add_text_box(slide, emu(1.6), emu(2.15), emu(2.0), emu(0.5), "SAM  ₩33.6B",
             font_size=12, bold=True, color=NAVY, font_name='Calibri')
add_text_box(slide, emu(0.9), emu(1.4), emu(2.5), emu(0.5), "TAM  ₩67.2B",
             font_size=12, bold=True, color=NAVY, font_name='Calibri')

# Right side - market detail cards
market_data = [
    ("TAM", "₩672억", "전체 수의 소프트웨어 시장",
     "전국 5,600 동물병원 × 월 10만원\nEMR 구독 기준 연간 시장 규모"),
    ("SAM", "₩336억", "수의 처방 보조 시장",
     "DUR 기능이 필요한 동물병원\n(다제 처방 활발한 중대형 병원 50%)"),
    ("SOM", "₩34억", "초기 진입 목표 시장",
     "EMR 파트너사 연동을 통해\n1~3년 내 확보 가능한 병원 500개\n× 월 5,000원 × 12개월\n+ 프리미엄 옵션 수익"),
]

for i, (label, amount, title, desc) in enumerate(market_data):
    x = 7.0
    y = 1.2 + i * 2.0

    card = add_card(slide, x, y, 5.8, 1.8, WHITE)

    # Label badge
    badge_color = [ICE_BLUE, MID_NAVY, NAVY][i]
    text_color = [NAVY, WHITE, WHITE][i]
    badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + 0.2), emu(y + 0.2), emu(0.75), emu(0.4), fill_color=badge_color)
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # Amount
    add_text_box(slide, emu(x + 1.15), emu(y + 0.15), emu(2.0), emu(0.45), amount,
                 font_size=22, bold=True, color=NAVY, font_name='Calibri')
    # Title
    add_text_box(slide, emu(x + 3.2), emu(y + 0.2), emu(2.3), emu(0.4), title,
                 font_size=11, bold=True, color=CHARCOAL, font_name='Calibri')
    # Description
    add_text_box(slide, emu(x + 0.2), emu(y + 0.7), emu(5.3), emu(1.0), desc,
                 font_size=10, bold=False, color=CHARCOAL, font_name='Calibri')

# Bottom note
add_text_box(slide, emu(0.8), emu(6.95), emu(12), emu(0.4),
             "현재 국내 경쟁자 전무  |  글로벌 사업자(Instinct Science) 한국 진출 유인 부족  |  국내 펫테크는 B2C 집중",
             font_size=11, bold=False, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 6, 12)


# ============================================================
# SLIDE 8: BUSINESS MODEL & EMR STRATEGY (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "BUSINESS MODEL",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "SaaS 구독 + 데이터 라이선싱 2단계 수익 모델",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Phase 1 Card
p1_card = add_card(slide, 0.6, 1.3, 5.8, 3.3, WHITE)
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.6), emu(1.3), emu(5.8), emu(0.06), fill_color=NAVY)

add_text_box(slide, emu(0.9), emu(1.5), emu(2.0), emu(0.4), "PHASE 1",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.9), emu(1.85), emu(5.2), emu(0.4), "SaaS 구독 수익",
             font_size=18, bold=True, color=NAVY, font_name='Calibri')

# EMR strategy visual
emr_items = [
    "EMR 회사에 API 플러그인 무상 제공",
    "병원 구독료 ₩5,000/월 (선택 옵션)",
    "수익 분배: DUR 70% / EMR 30%",
    "DUR 실제 매출: ₩7,000/병원/월",
]
for j, item in enumerate(emr_items):
    y_pos = 2.4 + j * 0.45
    dot = add_shape(slide, MSO_SHAPE.OVAL, emu(1.1), emu(y_pos + 0.08), emu(0.12), emu(0.12), fill_color=NAVY)
    dot.line.fill.background()
    add_text_box(slide, emu(1.4), emu(y_pos - 0.03), emu(4.8), emu(0.4), item,
                 font_size=12, bold=False, color=CHARCOAL, font_name='Calibri')

# Phase 2 Card
p2_card = add_card(slide, 6.9, 1.3, 5.8, 3.3, WHITE)
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(6.9), emu(1.3), emu(5.8), emu(0.06), fill_color=ACCENT_GREEN)

add_text_box(slide, emu(7.2), emu(1.5), emu(2.0), emu(0.4), "PHASE 2",
             font_size=11, bold=True, color=ACCENT_GREEN, font_name='Calibri')
add_text_box(slide, emu(7.2), emu(1.85), emu(5.2), emu(0.4), "데이터 라이선싱",
             font_size=18, bold=True, color=NAVY, font_name='Calibri')

data_items = [
    "비식별 처방 데이터 축적",
    "제약사 신약 R&D 데이터 판매",
    "사료/보조제 회사 맞춤 데이터",
    "시장 데이터 리포트: 연 ₩3,000만",
]
for j, item in enumerate(data_items):
    y_pos = 2.4 + j * 0.45
    dot = add_shape(slide, MSO_SHAPE.OVAL, emu(7.4), emu(y_pos + 0.08), emu(0.12), emu(0.12), fill_color=ACCENT_GREEN)
    dot.line.fill.background()
    add_text_box(slide, emu(7.7), emu(y_pos - 0.03), emu(4.8), emu(0.4), item,
                 font_size=12, bold=False, color=CHARCOAL, font_name='Calibri')

# Revenue projections
rev_card = add_card(slide, 0.6, 4.85, 12.1, 2.3, WHITE)
add_text_box(slide, emu(0.9), emu(5.0), emu(4), emu(0.35), "예상 매출 추이 (연간)",
             font_size=14, bold=True, color=NAVY, font_name='Calibri')

# Revenue bar chart simulation
years = ["1년차", "2년차", "3년차", "4년차", "5년차"]
revenues = [8.4, 42, 126, 168, 252]  # in millions KRW
max_rev = 252
bar_max_width = 7.0

for i, (year, rev) in enumerate(zip(years, revenues)):
    y = 5.5 + i * 0.3
    add_text_box(slide, emu(0.9), emu(y - 0.05), emu(0.8), emu(0.3), year,
                 font_size=9, bold=True, color=CHARCOAL, font_name='Calibri')
    # Bar
    bar_w = max(0.2, (rev / max_rev) * bar_max_width)
    bar_color = [ICE_BLUE, ICE_BLUE, SOFT_NAVY, MID_NAVY, NAVY][i]
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(1.8), emu(y), emu(bar_w), emu(0.2), fill_color=bar_color)
    # Amount label
    add_text_box(slide, emu(1.8 + bar_w + 0.1), emu(y - 0.05), emu(1.5), emu(0.3), f"₩{rev:.0f}M",
                 font_size=9, bold=True, color=NAVY, font_name='Calibri')

# Key metric callouts on right
for i, (num, label) in enumerate([("₩252M", "5년차 연매출"), ("3,000", "목표 병원 수"), ("₩120M", "5년차 순이익")]):
    x = 10.0
    y = 5.2 + i * 0.65
    add_text_box(slide, emu(x), emu(y), emu(1.8), emu(0.35), num,
                 font_size=18, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, emu(x), emu(y + 0.3), emu(1.8), emu(0.25), label,
                 font_size=9, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 7, 12)


# ============================================================
# SLIDE 9: TRACTION & TIMELINE (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "TRACTION & ROADMAP",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "현재 진행 상황 및 사업 추진 일정",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Current traction stats - horizontal row
traction_stats = [
    ("641+", "약물 DB 구축 완료"),
    ("10", "안전 검증 규칙"),
    ("100%", "MVP 엔진 구현"),
    ("7", "데모 임상 시나리오"),
    ("28", "PK 모델링 약물"),
]
for i, (num, label) in enumerate(traction_stats):
    x = 0.5 + i * 2.55
    card = add_card(slide, x, 1.2, 2.3, 1.3, WHITE)
    add_text_box(slide, emu(x), emu(1.3), emu(2.3), emu(0.6), num,
                 font_size=30, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, emu(x), emu(1.85), emu(2.3), emu(0.4), label,
                 font_size=10, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Timeline
timeline_items = [
    ("2026.06", "개발 & 윤리 심의", "법인 세팅, 품질 책임자 선임\n검역본부 임상시험 허가 요청", True),
    ("2026.07", "인허가 신청", "의료기기 제조업 허가\n제조품목 허가 동시 신청", True),
    ("2026.08\n~2027.02", "출시 준비", "수의학 교수진 임상 검증 논문\nEMR 파트너사 API 베타 테스트", False),
    ("2027.03\n~04", "시장 진출", "허가증 발급 → 즉각 상용화\n크라우드펀딩 + 자문위원회 구성", False),
]

# Timeline line
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.8), emu(4.15), emu(11.7), emu(0.04), fill_color=NAVY)

for i, (date, title, desc, is_done) in enumerate(timeline_items):
    x = 0.8 + i * 3.1

    # Timeline node
    node_color = ACCENT_GREEN if is_done else ICE_BLUE
    node = add_shape(slide, MSO_SHAPE.OVAL, emu(x + 1.2), emu(3.95), emu(0.45), emu(0.45), fill_color=node_color)
    node.line.fill.background()
    if is_done:
        tf = node.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Date above
    add_text_box(slide, emu(x), emu(3.0), emu(2.8), emu(0.8), date,
                 font_size=11, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)

    # Card below
    card = add_card(slide, x, 4.6, 2.8, 2.1, WHITE)
    add_text_box(slide, emu(x + 0.2), emu(4.75), emu(2.4), emu(0.4), title,
                 font_size=13, bold=True, color=NAVY, font_name='Calibri')
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 0.2), emu(5.15), emu(1.2), emu(0.025), fill_color=ICE_BLUE)
    add_text_box(slide, emu(x + 0.2), emu(5.3), emu(2.4), emu(1.2), desc,
                 font_size=10, bold=False, color=CHARCOAL, font_name='Calibri')

add_nav_dots(slide, 8, 12)


# ============================================================
# SLIDE 10: FUNDING & FINANCIALS (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "FUNDING & FINANCIALS",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "자금 소요 및 조달 계획",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Left - Expense breakdown
exp_card = add_card(slide, 0.6, 1.2, 6.0, 3.4, WHITE)
add_text_box(slide, emu(0.9), emu(1.35), emu(5.5), emu(0.4), "예상 자금 소요 (3년)",
             font_size=15, bold=True, color=NAVY, font_name='Calibri')

# Expense table
exp_headers = ["구분", "1년차", "2년차", "3년차", "합계"]
exp_data = [
    ["인건비", "18M", "36M", "54M", "108M"],
    ["서버/클라우드", "3M", "6M", "10M", "19M"],
    ["운영비/마케팅", "12M", "12M", "12M", "36M"],
    ["임대료", "12M", "12M", "12M", "36M"],
    ["합계", "45M", "66M", "88M", "199M"],
]

for j, h in enumerate(exp_headers):
    x = 0.9 + j * 1.1
    add_text_box(slide, emu(x), emu(1.85), emu(1.0), emu(0.3), h,
                 font_size=9, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)
# Header bg
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.9), emu(1.85), emu(5.5), emu(0.3), fill_color=NAVY)
for j, h in enumerate(exp_headers):
    x = 0.9 + j * 1.1
    add_text_box(slide, emu(x), emu(1.85), emu(1.0), emu(0.3), h,
                 font_size=9, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

for row_i, row in enumerate(exp_data):
    y = 2.2 + row_i * 0.35
    is_total = row_i == len(exp_data) - 1
    if is_total:
        add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.9), emu(y), emu(5.5), emu(0.35), fill_color=LIGHT_ICE)
    for j, val in enumerate(row):
        x = 0.9 + j * 1.1
        add_text_box(slide, emu(x), emu(y), emu(1.0), emu(0.3), val,
                     font_size=9, bold=is_total, color=NAVY if is_total else CHARCOAL,
                     font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Right - Funding sources
fund_card = add_card(slide, 6.9, 1.2, 5.8, 3.4, WHITE)
add_text_box(slide, emu(7.2), emu(1.35), emu(5.2), emu(0.4), "자금 조달 계획",
             font_size=15, bold=True, color=NAVY, font_name='Calibri')

funding_sources = [
    ("창업자 자금", "₩50M", "Pre-seed, Founder Capital", ICE_BLUE),
    ("엔젤 투자", "₩100M", "Angel Funding", SOFT_NAVY),
    ("정부지원금", "₩50M", "R&D Grant", MID_NAVY),
    ("창업경연대회", "N.A", "Startup Prize Money", NAVY),
]

for i, (name, amount, sub, color) in enumerate(funding_sources):
    y = 1.9 + i * 0.6
    # Color bar
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(7.4), emu(y + 0.05), emu(0.08), emu(0.4), fill_color=color)
    add_text_box(slide, emu(7.7), emu(y), emu(2.5), emu(0.3), name,
                 font_size=11, bold=True, color=NAVY, font_name='Calibri')
    add_text_box(slide, emu(7.7), emu(y + 0.25), emu(2.5), emu(0.25), sub,
                 font_size=8, bold=False, color=CHARCOAL, font_name='Calibri')
    add_text_box(slide, emu(10.5), emu(y + 0.05), emu(1.8), emu(0.35), amount,
                 font_size=14, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.RIGHT)

# Total
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(7.4), emu(4.3), emu(5.0), emu(0.03), fill_color=NAVY)
add_text_box(slide, emu(7.7), emu(4.35), emu(2.5), emu(0.3), "합계",
             font_size=12, bold=True, color=NAVY, font_name='Calibri')
add_text_box(slide, emu(10.5), emu(4.35), emu(1.8), emu(0.3), "₩200M+",
             font_size=16, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.RIGHT)

# Bottom - P&L Summary
pnl_card = add_card(slide, 0.6, 4.85, 12.1, 2.3, WHITE)
add_text_box(slide, emu(0.9), emu(5.0), emu(4), emu(0.35), "연도별 손익 추정",
             font_size=14, bold=True, color=NAVY, font_name='Calibri')

pnl_data = [
    ("", "1년차", "2년차", "3년차", "4년차", "5년차"),
    ("연 매출", "₩8.4M", "₩42M", "₩126M", "₩168M", "₩252M"),
    ("연 비용", "₩45M", "₩66M", "₩88M", "₩111M", "₩132M"),
    ("손익", "-₩36.6M", "-₩24M", "+₩38M", "+₩57M", "+₩120M"),
]

for row_i, row in enumerate(pnl_data):
    y = 5.45 + row_i * 0.35
    is_header = row_i == 0
    is_profit = row_i == 3

    if is_header:
        add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.9), emu(y), emu(11.5), emu(0.35), fill_color=NAVY)
    if is_profit:
        add_shape(slide, MSO_SHAPE.RECTANGLE, emu(0.9), emu(y), emu(11.5), emu(0.35), fill_color=LIGHT_ICE)

    for j, val in enumerate(row):
        x = 0.9 + j * 1.9
        text_color = WHITE if is_header else (ACCENT_GREEN if is_profit and '+' in val else (ACCENT_RED if is_profit and '-' in val else CHARCOAL))
        add_text_box(slide, emu(x), emu(y + 0.02), emu(1.8), emu(0.3), val,
                     font_size=10, bold=is_header or is_profit, color=text_color,
                     font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 9, 12)


# ============================================================
# SLIDE 11: TEAM (Light BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
light_bg(slide)

header_strip = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, emu(0.9), fill_color=NAVY)
header_strip.line.fill.background()
add_text_box(slide, emu(0.8), emu(0.2), emu(8), emu(0.5), "OUR TEAM",
             font_size=11, bold=True, color=ICE_BLUE, font_name='Calibri')
add_text_box(slide, emu(0.8), emu(0.42), emu(10), emu(0.45), "창업팀 소개",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')

# Team member cards
team_members = [
    ("염인", "CEO", "카네기멜론대학교\nHCI / Computer Science",
     "팀총괄, 사업 모델 설계\n시스템 아키텍처 설계\n해외 파트너쉽, 펀딩", NAVY),
    ("김명학", "CMO", "OO대학교\n법규/행정 전문",
     "재무/회계, 마케팅\n시장조사/분석\n국내 파트너쉽, 영업총괄", MID_NAVY),
    ("김동현", "CTO", "세종대학교\n인공지능 데이터 사이언스",
     "제품 알고리즘 개발\n약물 데이터베이스 구축\nAPI/EMR 연동 시스템", SOFT_NAVY),
    ("김준우", "QM", "울산과기대 전자공학과\nLG시스플레이 선임연구원",
     "의료기기법 품질관리\n제6조 품질책임자\n인허가 프로세스", RGBColor(0x4A, 0x5A, 0xB8)),
]

for i, (name, role, edu, responsibilities, color) in enumerate(team_members):
    x = 0.5 + i * 3.2

    # Card
    card = add_card(slide, x, 1.2, 2.95, 5.4, WHITE)

    # Avatar circle
    avatar = add_shape(slide, MSO_SHAPE.OVAL, emu(x + 0.85), emu(1.4), emu(1.2), emu(1.2), fill_color=color)
    avatar.line.fill.background()
    tf = avatar.text_frame
    p = tf.paragraphs[0]
    p.text = name[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Name + Role
    add_text_box(slide, emu(x + 0.2), emu(2.8), emu(2.55), emu(0.4), name,
                 font_size=18, bold=True, color=NAVY, font_name='Calibri', alignment=PP_ALIGN.CENTER)

    # Role badge
    role_badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + 0.85), emu(3.25), emu(1.2), emu(0.35), fill_color=color)
    role_badge.line.fill.background()
    tf = role_badge.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Education
    add_text_box(slide, emu(x + 0.2), emu(3.8), emu(2.55), emu(0.7), edu,
                 font_size=10, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, emu(x + 0.5), emu(4.5), emu(1.95), emu(0.025), fill_color=ICE_BLUE)

    # Responsibilities
    add_text_box(slide, emu(x + 0.2), emu(4.7), emu(2.55), emu(1.8), responsibilities,
                 font_size=10, bold=False, color=CHARCOAL, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Bottom affiliation bar
affil_bar = add_card(slide, 0.6, 6.8, 12.1, 0.5, NAVY)
add_text_box(slide, emu(0.8), emu(6.85), emu(12), emu(0.4),
             "제75사단 포병여단 527 포병대대  |  카네기멜론대학교  |  세종대학교 TEED LAB  |  울산과기대",
             font_size=10, bold=False, color=ICE_BLUE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 10, 12)


# ============================================================
# SLIDE 12: CONTACT / CLOSING (Dark BG)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)

# Decorative elements (mirror of cover)
# Dot grid
for row in range(12):
    for col in range(20):
        x = 0.5 + col * 0.6
        y = 0.3 + row * 0.6
        dot = add_shape(slide, MSO_SHAPE.OVAL, emu(x), emu(y), emu(0.04), emu(0.04), fill_color=RGBColor(0x28, 0x34, 0x75))
        dot.line.fill.background()

# Large decorative circles
for cx, cy, sz in [(1.0, 5.0, 2.5), (11.0, 0.5, 1.8), (10.0, 5.5, 1.2)]:
    c = add_shape(slide, MSO_SHAPE.OVAL, emu(cx), emu(cy), emu(sz), emu(sz), fill_color=RGBColor(0x22, 0x2E, 0x6A))
    c.line.fill.background()

# Center content
# Logo
add_text_box(slide, emu(0), emu(1.5), emu(13.333), emu(0.5), "NUVOVET",
             font_size=18, bold=True, color=ICE_BLUE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, emu(2), emu(2.2), emu(9.333), emu(1.0),
             "동물 의료의 처방 안전망을\n함께 만들어 갑니다",
             font_size=36, bold=True, color=WHITE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Accent line
add_shape(slide, MSO_SHAPE.RECTANGLE, emu(5.5), emu(3.5), emu(2.333), emu(0.05), fill_color=ICE_BLUE)

# Contact cards
contact_card = add_card(slide, 3.5, 4.0, 6.333, 2.2, RGBColor(0x1A, 0x23, 0x55))

add_text_box(slide, emu(3.8), emu(4.2), emu(5.7), emu(0.35), "CONTACT US",
             font_size=12, bold=True, color=ICE_BLUE, font_name='Calibri', alignment=PP_ALIGN.CENTER)

contacts = [
    "CEO 염인  |  카네기멜론대학교 HCI/CS",
    "CTO 김동현  |  세종대학교 AI 데이터사이언스",
    "github.com/nuvovet_systems",
]
for i, c in enumerate(contacts):
    add_text_box(slide, emu(3.8), emu(4.7 + i * 0.4), emu(5.7), emu(0.35), c,
                 font_size=12, bold=False, color=RGBColor(0xAA, 0xBB, 0xDD), font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, emu(0), emu(6.5), emu(13.333), emu(0.4),
             "말 못하는 생명을 위한, 데이터 기반 처방 안전망",
             font_size=14, bold=False, color=RGBColor(0x77, 0x88, 0xAA), font_name='Calibri', alignment=PP_ALIGN.CENTER)

add_nav_dots(slide, 11, 12)


# ============================================================
# SAVE
# ============================================================
output_path = "/mnt/user-data/outputs/nuvovet_pitch_deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
